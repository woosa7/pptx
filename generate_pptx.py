import sys
import os
import datetime
from pathlib import Path
from PyQt5.QtWidgets import *
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor

save_filename = ''

class MyWindow(QMainWindow):
    def __init__(self):
        super().__init__()
        self.setupUI()

    def setupUI(self):
        self.setGeometry(800, 400, 300, 120)

        # Label
        label = QLabel("Name of file:", self)
        label.move(20, 20)

        # LineEdit
        self.lineEdit = QLineEdit("", self)
        self.lineEdit.setGeometry(100,25,120,20)
        self.lineEdit.textChanged.connect(self.lineEditChanged)

        # Label
        label2 = QLabel(".pptx", self)
        label2.move(223, 20)

        # Button
        btn1 = QPushButton("OK", self)
        btn1.move(20, 60)
        btn1.clicked.connect(self.btn1_clicked)

        btn2 = QPushButton("Clear", self)
        btn2.move(150, 60)
        btn2.clicked.connect(self.btn2_clicked)

        # StatusBar
        self.statusBar = QStatusBar(self)
        self.setStatusBar(self.statusBar)

    def lineEditChanged(self):
        self.statusBar.showMessage(self.lineEdit.text()+'.pptx')

    def btn1_clicked(self):
        save_filename = self.lineEdit.text()
        if len(save_filename) == 0:
            self.statusBar.showMessage('Enter the name of file!')
        else:
            self.close()

    def btn2_clicked(self):
        self.lineEdit.clear()


app = QApplication(sys.argv)
mywindow = MyWindow()
mywindow.show()
app.exec_()

save_filename = mywindow.lineEdit.text()
if len(save_filename) == 0:
    save_filename = 'noname'

# Test
# save_filename = '111'

save_filename = save_filename+'.pptx'

# -----------------------------------------------------------------------------------------
# predefine parameters

main_title = ""
writer = ""
present_date = ""
titlebar_background = ""
department1 = "Nano/Bio Computational Chemistry Lab."
department2 = "Department of Chemistry, Sookmyung Women’s University, Seoul, Korea"

# -----------------------------------------------------------------------------------------
# Read parameter file
title_list = []
fig_list = []
img_list = []
desc_list = []

params_file = './ppt_contents.txt'
if Path(params_file).is_file():
    for line in open(params_file, 'r'):
        str = line.strip()

        idx = str.find('"')
        cont = str[idx+1:].replace('"', '')
        if len(cont) > 0 and str[:1] != '#':
            if str[:4] == 'page':
                if 'title' in str:
                    title_list.append(str)
                if 'desc' in str:
                    desc_list.append(str)
                if 'fig ' in str:
                    fig_list.append(str)
                if 'fig_' in str:
                    img_list.append(str)
            else:
                exec(str)


# -----------------------------------------------------------------------------------------
# Main Title Page

prs = Presentation('./default.pptx')

blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

# add_(left, top, width, height) : max width = 25.4

# logo
image_file1 = './images/nbcc_logo1.png'
image_file2 = './images/nbcc_logo2.png'
slide.shapes.add_picture(image_file1, Cm(0), Cm(0), Cm(2.2))
slide.shapes.add_picture(image_file2, Cm(23.2), Cm(0), Cm(2.2))

# Main Title
text_box = slide.shapes.add_textbox(Cm(1.0), Cm(4.0), Cm(23.4), Cm(3.0))
text_box.text_frame.word_wrap = True
p = text_box.text_frame.paragraphs[0]
p.text = main_title
p.font.bold = True
p.font.size = Pt(32)
p.font.name = 'Arial Black'
p.alignment = PP_ALIGN.CENTER

# Footer
text_box = slide.shapes.add_textbox(Cm(0.0), Cm(17.0), Cm(25.4), Cm(1.6))
p1 = text_box.text_frame.paragraphs[0]
p1.text = 'Nano/Bio Computational Chemistry Lab.'
p1.font.bold = True
p1.font.italic = True
p1.font.size = Pt(14)
p1.font.name = 'Malgun Gothic'
p1.alignment = PP_ALIGN.CENTER

p2 = text_box.text_frame.add_paragraph()
p2.text = 'Department of Chemistry, Sookmyung Women’s University, Seoul, Korea'
p2.font.size = Pt(14)
p2.font.name = 'Malgun Gothic'
p2.alignment = PP_ALIGN.CENTER

# Writer & Date
if len(writer) == 0:
    writer ="Your Name"
if len(present_date) == 0:
    present_date = datetime.datetime.now().strftime('%Y.%m.%d.')

text_box = slide.shapes.add_textbox(Cm(1.0), Cm(14.0), Cm(23.4), Cm(1.5))
p = text_box.text_frame.paragraphs[0]
p.text = writer + '\n' + present_date
p.font.size = Pt(14)
p.font.name = 'Malgun Gothic'
p.alignment = PP_ALIGN.CENTER

#-----------------------------------------------------------------------------------------
# Sub-title pages

page_len = len(title_list)

for k in range(page_len):
    slide = prs.slides.add_slide(blank_slide_layout)

    # Title background image
    exist_titlebar = False
    if Path(titlebar_background).is_file():
        slide.shapes.add_picture(titlebar_background, Cm(0), Cm(0), Cm(25.4), Cm(2.28))
        exist_titlebar = True
    else:
        # 이미 없으면 타이틀 아래 밑줄
        shapes = slide.shapes
        shape = shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1.0), Cm(2.2), Cm(23.4), Cm(0.22)
        )
        fill = shape.fill
        fill.solid()
        fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
        fill.fore_color.brightness = -0.25

    # Title
    text_box = slide.shapes.add_textbox(Cm(1.0), Cm(0.3), Cm(23.4), Cm(1.62))
    p = text_box.text_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.name = 'Arial'
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    if exist_titlebar:
        p.font.color.rgb = RGBColor(255, 255, 255)
    else:
        p.font.color.rgb = RGBColor(0, 0, 0)

    content = title_list[k]
    if content[:5] == 'page{}'.format(k+1):
        idx = content.find('"')
        p.text = content[idx+1:].replace('"', '')
    elif content[:6] == 'page{}'.format(k + 1):
        idx = content.find('"')
        p.text = content[idx + 1:].replace('"', '')
    else:
        p.text = 'Title'

    # Figure
    if len(fig_list) >= k + 1:  # fig 존재
        text_box = slide.shapes.add_textbox(Cm(1.5), Cm(16.0), Cm(22.4), Cm(1.0))
        p = text_box.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.name = 'Tahoma'
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT

        content = fig_list[k]
        if content[:5] == 'page{}'.format(k+1):
            idx = content.find('"')
            p.text = content[idx+1:].replace('"', '')
        else:
            p.text = 'Figure 0.'

    # Figure image & Description
    content = desc_list[k]

    idx_desc = content.find('"')
    desc = content[idx_desc + 1:].replace('"', '')

    if len(img_list) < k+1:   # 이미지 없음
        # description
        text_box = slide.shapes.add_textbox(Cm(1.5), Cm(5.0), Cm(22.4), Cm(10.0))
        text_box.text_frame.word_wrap = True
        p = text_box.text_frame.paragraphs[0]
        p.font.size = Pt(16)
        p.font.name = 'Arial'
        p.alignment = PP_ALIGN.LEFT
        p.text = desc

    else:  # 이미지 있는 경우
        fig_path = img_list[k]

        idx = fig_path.find('"')
        image_file = fig_path[idx+1:].replace('"', '')

        tbox_left = Cm(1.5)
        tbox_width = Cm(22.4)
        if os.path.exists(image_file):
            slide.shapes.add_picture(image_file, Cm(1.5), Cm(5))
            tbox_left = Cm(10.5)
            tbox_width = Cm(13.4)

        # description
        text_box = slide.shapes.add_textbox(tbox_left, Cm(5.0), tbox_width, Cm(10.0))
        text_box.text_frame.word_wrap = True
        p = text_box.text_frame.paragraphs[0]
        p.font.size = Pt(16)
        p.font.name = 'Arial'
        p.alignment = PP_ALIGN.LEFT
        p.text = desc



    # Page no.
    text_box = slide.shapes.add_textbox(Cm(24), Cm(18.0), Cm(1), Cm(0.7))
    p = text_box.text_frame.paragraphs[0]
    p.font.size = Pt(10)
    p.font.name = 'Arial'
    p.alignment = PP_ALIGN.RIGHT
    p.text = '{}'.format(k+1)


#-----------------------------------------------------------------------------------------
prs.save(save_filename)

os.system(save_filename)

sys.exit()


# pyinstaller --noconsole --onefile generate_pptx.py
#-----------------------------------------------------------------------------------------